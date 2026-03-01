import io
from typing import List
from core.models import DocumentChunk, FileType
from .base_extractor import BaseExtractor

class PowerPointExtractor(BaseExtractor):
    SUPPORTED_EXTENSIONS = ["pptx", "ppt"]

    def extract(self, file_path: str) -> List[DocumentChunk]:
        chunks: List[DocumentChunk] = []
        source = file_path.split("/")[-1]
        chunk_idx = 0

        try:
            from pptx import Presentation
            from pptx.util import Inches
            import zipfile

            prs = Presentation(file_path)

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_texts.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            if row_text.strip():
                                slide_texts.append(row_text)

                if slide_texts:
                    chunks.append(self._make_text_chunk(
                        text=f"[Slide {slide_num}]\n" + "\n".join(slide_texts),
                        source=source,
                        page=slide_num,
                        chunk_index=chunk_idx,
                        file_type=FileType.POWERPOINT,
                    ))
                    chunk_idx += 1

            with zipfile.ZipFile(file_path, "r") as zf:
                image_files = [n for n in zf.namelist() if n.startswith("ppt/media/")]
                for img_path in image_files:
                    ext = img_path.rsplit(".", 1)[-1].lower()
                    if ext not in ("png", "jpg", "jpeg", "gif", "bmp", "tiff", "webp"):
                        continue
                    mime = f"image/{ext}" if ext != "jpg" else "image/jpeg"
                    image_data = zf.read(img_path)
                    chunks.append(self._make_image_chunk(
                        image_data=image_data,
                        mime=mime,
                        source=source,
                        page=0,
                        chunk_index=chunk_idx,
                        file_type=FileType.POWERPOINT,
                    ))
                    chunk_idx += 1

        except ImportError as e:
            raise RuntimeError(
                f"Missing dependency for PowerPoint extraction: {e}. "
                "Run: pip install python-pptx"
            )
        except Exception as e:
            raise RuntimeError(f"PowerPoint extraction failed for '{file_path}': {e}")

        return chunks